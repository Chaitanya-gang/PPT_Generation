"""
newd2p - Creative PPT Builder with Charts & Images
"""

import json
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.config import get_settings
from src.utils.logger import get_logger
from src.images.image_generator import generate_image_from_prompt
from src.diagrams.diagram_generator import create_flow_diagram

logger = get_logger("ppt_builder")
settings = get_settings()


from src.ppt.theme_manager import get_theme

TOPIC_ICONS = {
    "introduction": "🌟",
    "overview": "🔭",
    "problem": "⚠️",
    "solution": "💡",
    "data": "📊",
    "statistics": "📈",
    "result": "🎯",
    "results": "🎯",
    "impact": "💥",
    "environment": "🌍",
    "technology": "🚀",
    "health": "❤️",
    "education": "📚",
    "finance": "💰",
    "conclusion": "✅",
    "future": "🔮",
    "action": "⚡",
    "climate": "🌡️",
    "food": "🌾",
    "water": "💧",
    "energy": "⚡",
    "state": "📋",
    "human": "👥",
    "key": "🔑",
    "proposed": "🛠️",
    "default": "✦",
}


def get_icon(title: str) -> str:
    title_lower = title.lower()
    for key, icon in TOPIC_ICONS.items():
        if key in title_lower:
            return icon
    return TOPIC_ICONS["default"]


class PPTBuilder:

    def __init__(
        self,
        theme_name: str = "vibrant",
        image_mode: bool = False,
        diagram_mode: bool = False,
        include_speaker_notes: bool = True,
    ):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.colors = get_theme(theme_name)
        if not self.colors:
            # Emergency fallback if no themes loaded
            self.colors = {
                "bg": RGBColor(15, 23, 42),
                "card_bg": RGBColor(30, 41, 59),
                "title": RGBColor(255, 255, 255),
                "text": RGBColor(226, 232, 240),
                "accent1": RGBColor(99, 102, 241),
                "accent2": RGBColor(244, 63, 94),
                "accent3": RGBColor(34, 211, 238),
                "accent4": RGBColor(163, 230, 53),
                "highlight": RGBColor(251, 191, 36),
                "muted": RGBColor(148, 163, 184),
            }
        self.theme_name = theme_name
        self.total_slides = 0
        self.chart_paths = []
        self.diagram_paths = []
        self.image_mode = image_mode
        self.diagram_mode = diagram_mode
        self.include_speaker_notes = include_speaker_notes

    def set_charts(self, chart_paths: list):
        self.chart_paths = chart_paths

    def set_diagrams(self, diagram_paths: list):
        """
        Attach pre-generated diagrams (from automatic analysis) to be used in slides.
        Each item should be a dict with at least {"path": ..., "title": ...}.
        """
        self.diagram_paths = diagram_paths

    def build_from_json(self, narrative_json: str, output_path: str) -> str:
        try:
            data = json.loads(narrative_json)
        except json.JSONDecodeError:
            data = self._extract_json(narrative_json)

        title = data.get("title", "Presentation")
        subtitle = data.get("subtitle", "")
        slides = data.get("slides", [])
        # Total slides includes potential extra chart-only slides
        self.total_slides = len(slides) + len(self.chart_paths)

        logger.info(f"Building creative PPT: '{title}' with {len(slides)} slides")

        chart_index = 0
        diagram_index = 0
        current_number = 0

        for slide_data in slides:
            current_number += 1
            visual_cue = slide_data.get("visual_cue", "")
            slide_number = slide_data.get("slide_number", 0)

            # Optional auto image generation
            if self.image_mode and visual_cue:
                prompt = f"{title} - {slide_data.get('title', '')}: {visual_cue}"
                img_output = Path(settings.image_output_dir) / f"{slide_number:03d}_image.png"
                image_path = generate_image_from_prompt(prompt, str(img_output))
                if image_path:
                    slide_data["image_path"] = image_path

            # Optional simple flow diagram (legacy rule-based + infographic keywords)
            visual_lower = visual_cue.lower()
            if self.diagram_mode and any(
                kw in visual_lower for kw in ["diagram", "infographic", "workflow", "process", "flow"]
            ):
                stages = []
                if "input" in visual_lower:
                    stages.append("Input")
                if "process" in visual_lower or "model" in visual_lower:
                    stages.append("Model / Process")
                if "output" in visual_lower:
                    stages.append("Output")
                if not stages:
                    stages = ["Input", "Process", "Output"]

                diag_output = Path(settings.diagram_output_dir) / f"{slide_number:03d}_diagram"
                diagram_path = create_flow_diagram(
                    title=slide_data.get("title", ""),
                    stages=stages,
                    output_path=str(diag_output),
                )
                if diagram_path:
                    slide_data["diagram_path"] = diagram_path

            # Attach any pre-generated auto diagrams to diagram slides
            slide_type = slide_data.get("slide_type", "content")
            layout = slide_data.get("layout", "").lower()

            # Ensure consistent slide_number for badges
            slide_data["slide_number"] = current_number

            if slide_type == "title":
                self._add_title_slide(slide_data, subtitle)
            elif layout == "title_big_number" or slide_type == "big_number":
                self._add_big_number_slide(slide_data)
            elif slide_type == "content":
                self._add_content_slide(slide_data)
            elif slide_type == "transition":
                self._add_transition_slide(slide_data)
            elif slide_type == "chart":
                chart_path = None
                if chart_index < len(self.chart_paths):
                    chart_path = self.chart_paths[chart_index].get("path")
                    chart_index += 1
                self._add_chart_slide(slide_data, chart_path)
            elif slide_type == "closing":
                self._add_closing_slide(slide_data)
            else:
                self._add_content_slide(slide_data)

        # Add remaining charts as extra slides
        while chart_index < len(self.chart_paths):
            chart_info = self.chart_paths[chart_index]
            current_number += 1
            self._add_chart_slide(
                {"title": chart_info["title"], "slide_number": current_number},
                chart_info["path"],
            )
            chart_index += 1

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        logger.info(f"PPT saved: {output_path}")
        return output_path

    def _set_bg(self, slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["bg"]

    def _add_shape(self, slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _add_rounded_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _add_text(self, slide, left, top, width, height, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_pill(self, slide, left, top, width, height, text, fill_color, text_color):
        pill = self._add_rounded_rect(slide, left, top, width, height, fill_color)
        tf = pill.text_frame
        tf.word_wrap = False
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = text_color
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        return pill

    def _add_stat_card(self, slide, left, top, width, height, number, label, color):
        """Add a statistical highlight card"""
        self._add_rounded_rect(slide, left, top, width, height, self.colors["card_bg"])
        self._add_shape(slide, left, top, width, Inches(0.06), color)

        self._add_text(slide,
            left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.8),
            number, Pt(28), color, bold=True, alignment=PP_ALIGN.CENTER
        )
        self._add_text(slide,
            left + Inches(0.2), top + Inches(0.9), width - Inches(0.4), Inches(0.5),
            label, Pt(11), self.colors["muted"], alignment=PP_ALIGN.CENTER
        )

    def _add_title_slide(self, data: dict, subtitle: str = ""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide)

        self._add_shape(slide, Inches(9), Inches(-2), Inches(6), Inches(6),
            self.colors["accent1"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(10.5), Inches(-0.5), Inches(4), Inches(4),
            self.colors["accent2"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(-1), Inches(5), Inches(3), Inches(3),
            self.colors["accent3"], MSO_SHAPE.OVAL)

        self._add_shape(slide, Inches(1), Inches(2.8), Inches(2), Inches(0.08),
            self.colors["highlight"])

        self._add_pill(
            slide,
            Inches(1),
            Inches(0.55),
            Inches(2.3),
            Inches(0.45),
            self.theme_name.upper(),
            self.colors["card_bg"],
            self.colors["accent3"],
        )

        title = data.get("title", "")
        self._add_text(slide, Inches(1), Inches(1.2), Inches(9), Inches(1.8),
            title, Pt(44), self.colors["title"], bold=True)

        sub = subtitle or data.get("subtitle", "")
        if sub:
            self._add_text(slide, Inches(1), Inches(3.2), Inches(9), Inches(1),
                sub, Pt(20), self.colors["muted"])

        self._add_rounded_rect(slide, Inches(1), Inches(5), Inches(3.5), Inches(1.5),
            self.colors["card_bg"])
        self._add_text(slide, Inches(1.3), Inches(5.1), Inches(3), Inches(0.5),
            "📊 Data-Driven Insights", Pt(13), self.colors["accent3"], bold=True)
        self._add_text(slide, Inches(1.3), Inches(5.6), Inches(3), Inches(0.7),
            "Powered by AI Analysis", Pt(11), self.colors["muted"])

        self._add_rounded_rect(slide, Inches(5), Inches(5), Inches(3.5), Inches(1.5),
            self.colors["card_bg"])
        self._add_text(slide, Inches(5.3), Inches(5.1), Inches(3), Inches(0.5),
            "🎯 Actionable Results", Pt(13), self.colors["accent2"], bold=True)
        self._add_text(slide, Inches(5.3), Inches(5.6), Inches(3), Inches(0.7),
            "Clear recommendations", Pt(11), self.colors["muted"])

        self._add_text(slide, Inches(1), Inches(6.8), Inches(5), Inches(0.5),
            "✦  Powered by newd2p AI", Pt(11), self.colors["muted"])

        self._add_rounded_rect(slide, Inches(9.2), Inches(4.8), Inches(3.1), Inches(1.8),
            self.colors["card_bg"])
        self._add_text(slide, Inches(9.45), Inches(5.05), Inches(2.7), Inches(0.4),
            "Presentation Flow", Pt(12), self.colors["accent4"], bold=True)
        self._add_text(slide, Inches(9.45), Inches(5.45), Inches(2.4), Inches(0.8),
            "Upload\nAnalyze\nGenerate", Pt(16), self.colors["text"], bold=True)

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_content_slide(self, data: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide)

        slide_num = data.get("slide_number", 0)
        title = data.get("title", "")
        bullets = data.get("bullet_points", [])
        icon = get_icon(title)
        image_path = data.get("image_path")
        diagram_path = data.get("diagram_path")
        visual = data.get("visual_cue", "")
        layout = data.get("layout", "title_bullets").lower()

        self._add_shape(slide, Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.06), self.colors["accent1"])

        badge = self._add_rounded_rect(slide,
            Inches(11.5), Inches(0.3), Inches(1.2), Inches(0.6),
            self.colors["accent1"])
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{slide_num} / {self.total_slides}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        self._add_text(slide, Inches(0.8), Inches(0.4), Inches(1), Inches(0.8),
            icon, Pt(36), self.colors["accent1"])

        self._add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(1),
            title, Pt(32), self.colors["title"], bold=True)

        self._add_shape(slide, Inches(0.8), Inches(2.0), Inches(4), Inches(0.04),
            self.colors["accent2"])

        self._add_pill(
            slide,
            Inches(0.8),
            Inches(0.34),
            Inches(2.1),
            Inches(0.4),
            layout.replace("_", " ").upper()[:18],
            self.colors["card_bg"],
            self.colors["muted"],
        )

        has_visual_asset = bool(
            (diagram_path and Path(diagram_path).exists()) or
            (image_path and Path(image_path).exists())
        )

        if layout == "two_column":
            self._render_bullets(slide, bullets, Inches(0.8), Inches(2.45), Inches(6.25), Inches(0.9))
            visual_x = Inches(8.0)
            visual_y = Inches(2.35)
            visual_w = Inches(4.0)
            visual_h = Inches(3.95)
        elif layout == "visual_focus":
            self._render_visual_direction_card(slide, visual, Inches(0.8), Inches(2.45), Inches(4.05), Inches(2.5))
            self._render_bullets(slide, bullets, Inches(5.2), Inches(2.45), Inches(6.9), Inches(0.9), compact=True)
            visual_x = Inches(0.95)
            visual_y = Inches(5.2)
            visual_w = Inches(3.7)
            visual_h = Inches(1.0)
        else:
            card_width = Inches(6.6) if has_visual_asset else Inches(11.5)
            self._render_bullets(slide, bullets, Inches(0.8), Inches(2.5), card_width, Inches(0.95))
            visual_x = Inches(8.2)
            visual_y = Inches(2.5)
            visual_w = Inches(3.8)
            visual_h = Inches(3.8)

        visual_asset_path: Optional[str] = None
        if diagram_path:
            visual_asset_path = diagram_path
        elif image_path:
            visual_asset_path = image_path

        if visual_asset_path and Path(visual_asset_path).exists():
            self._add_rounded_rect(
                slide,
                visual_x - Inches(0.2),
                visual_y - Inches(0.2),
                visual_w + Inches(0.4),
                visual_h + Inches(0.4),
                self.colors["card_bg"],
            )
            slide.shapes.add_picture(
                visual_asset_path,
                visual_x,
                visual_y,
                width=visual_w,
                height=visual_h,
            )
            self._add_pill(
                slide,
                visual_x - Inches(0.1),
                visual_y - Inches(0.55),
                Inches(1.7),
                Inches(0.38),
                "VISUAL",
                self.colors["card_bg"],
                self.colors["accent3"],
            )
        elif visual:
            self._render_visual_direction_card(slide, visual, visual_x - Inches(0.1), visual_y, visual_w, Inches(2.2))

        if visual:
            self._add_rounded_rect(slide,
                Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6),
                self.colors["card_bg"])
            self._add_text(slide,
                Inches(1.0), Inches(6.65), Inches(11), Inches(0.5),
                f"🎨  {visual}", Pt(10), self.colors["muted"])

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _render_bullets(self, slide, bullets, left, top, card_width, card_height, compact: bool = False):
        if not bullets:
            return

        accent_colors = [
            self.colors["accent1"], self.colors["accent2"],
            self.colors["accent3"], self.colors["accent4"],
        ]
        font_size = Pt(14 if compact else 16)
        vertical_gap = 0.9 if compact else 1.15

        for i, bullet in enumerate(bullets):
            y_pos = top + Inches(i * vertical_gap)
            accent = accent_colors[i % len(accent_colors)]

            self._add_rounded_rect(slide, left, y_pos, card_width, card_height, self.colors["card_bg"])
            self._add_shape(slide, left, y_pos, Inches(0.12), card_height, accent)

            num_circle = self._add_shape(
                slide,
                left + Inches(0.4),
                y_pos + Inches(0.15),
                Inches(0.55),
                Inches(0.55),
                accent,
                MSO_SHAPE.OVAL,
            )
            tf = num_circle.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            self._add_text(
                slide,
                left + Inches(1.2),
                y_pos + Inches(0.15),
                card_width - Inches(1.45),
                Inches(0.7),
                bullet,
                font_size,
                self.colors["text"],
            )

    def _render_visual_direction_card(self, slide, visual, left, top, width, height):
        self._add_rounded_rect(slide, left, top, width, height, self.colors["card_bg"])
        self._add_text(
            slide,
            left + Inches(0.2),
            top + Inches(0.2),
            width - Inches(0.4),
            Inches(0.4),
            "Visual Direction",
            Pt(12),
            self.colors["accent3"],
            bold=True,
        )
        self._add_text(
            slide,
            left + Inches(0.2),
            top + Inches(0.65),
            width - Inches(0.4),
            height - Inches(0.8),
            (visual or "Use a supporting visual here")[:180],
            Pt(15),
            self.colors["text"],
        )

    def _add_chart_slide(self, data: dict, chart_path: str = None):
        """Slide with embedded chart image - enhanced with modern design"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide)

        title = data.get("title", "Data Analysis")
        slide_num = data.get("slide_number", 0)
        icon = "📊"

        # Enhanced header
        self._add_shape(slide, Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.08), self.colors["accent3"])

        # Icon with background
        icon_bg = self._add_shape(slide,
            Inches(0.6), Inches(0.3), Inches(0.8), Inches(0.8),
            self.colors["accent3"], MSO_SHAPE.OVAL)
        self._add_text(slide, Inches(0.8), Inches(0.4), Inches(0.4), Inches(0.6),
            icon, Pt(24), RGBColor(255, 255, 255), alignment=PP_ALIGN.CENTER)

        self._add_text(slide, Inches(1.6), Inches(0.4), Inches(9.5), Inches(0.8),
            title, Pt(28), self.colors["title"], bold=True)

        # Decorative elements
        self._add_shape(slide, Inches(1.6), Inches(1.2), Inches(3), Inches(0.06),
            self.colors["accent3"])
        self._add_shape(slide, Inches(11), Inches(1.5), Inches(1), Inches(0.04),
            self.colors["highlight"])

        if chart_path and Path(chart_path).exists():
            # Enhanced chart container with shadow effect
            self._add_rounded_rect(slide,
                Inches(1.3), Inches(2.0), Inches(10.7), Inches(5.2),
                self.colors["card_bg"])
            self._add_rounded_rect(slide,
                Inches(1.4), Inches(2.1), Inches(10.5), Inches(5.0),
                RGBColor(
                    min(255, self.colors["card_bg"].rgb[0] + 15),
                    min(255, self.colors["card_bg"].rgb[1] + 15),
                    min(255, self.colors["card_bg"].rgb[2] + 15)
                ))
            slide.shapes.add_picture(
                chart_path,
                Inches(1.6), Inches(2.3),
                Inches(10.1), Inches(4.6)
            )

            # Add corner accent
            self._add_shape(slide, Inches(1.3), Inches(2.0), Inches(0.3), Inches(0.3),
                self.colors["accent3"], MSO_SHAPE.TRIANGLE)
        else:
            # Enhanced placeholder
            self._add_rounded_rect(slide,
                Inches(2), Inches(2.5), Inches(9), Inches(4),
                self.colors["card_bg"])
            self._add_shape(slide, Inches(2), Inches(2.5), Inches(9), Inches(0.08),
                self.colors["accent3"])  # Top accent

            self._add_text(slide,
                Inches(3), Inches(3.5), Inches(7), Inches(2),
                "📊 Chart will be generated from data analysis",
                Pt(20), self.colors["muted"], alignment=PP_ALIGN.CENTER)

        # Enhanced badge
        badge = self._add_rounded_rect(slide,
            Inches(11.5), Inches(0.3), Inches(1.2), Inches(0.6),
            self.colors["accent3"])
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{slide_num}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Add decorative elements
        self._add_shape(slide, Inches(10), Inches(6.5), Inches(2), Inches(0.08),
            self.colors["highlight"])
        self._add_shape(slide, Inches(0.5), Inches(6.8), Inches(1.5), Inches(0.04),
            self.colors["accent4"])

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_big_number_slide(self, data: dict):
        """Visual highlight slide with a single big metric."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide)

        title = data.get("title", "Key Result")
        metric = data.get("metric", "")
        description = data.get("description", "")
        slide_num = data.get("slide_number", 0)

        # Header stripe
        self._add_shape(
            slide,
            Inches(0),
            Inches(0),
            self.prs.slide_width,
            Inches(0.06),
            self.colors["accent4"],
        )

        # Slide counter badge
        badge = self._add_rounded_rect(
            slide,
            Inches(11.5),
            Inches(0.3),
            Inches(1.2),
            Inches(0.6),
            self.colors["accent4"],
        )
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num} / {self.total_slides}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        self._add_text(
            slide,
            Inches(0.8),
            Inches(0.9),
            Inches(11),
            Inches(1),
            title,
            Pt(32),
            self.colors["title"],
            bold=True,
        )

        # Big number
        self._add_text(
            slide,
            Inches(1.5),
            Inches(2.3),
            Inches(10.5),
            Inches(2.5),
            metric,
            Pt(80),
            self.colors["accent1"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Description
        if description:
            self._add_text(
                slide,
                Inches(2.0),
                Inches(4.4),
                Inches(9.5),
                Inches(1.0),
                description,
                Pt(22),
                self.colors["text"],
                alignment=PP_ALIGN.CENTER,
            )

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_transition_slide(self, data: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide)

        title = data.get("title", "")
        icon = get_icon(title)

        self._add_shape(slide, Inches(-2), Inches(-2), Inches(7), Inches(7),
            self.colors["accent1"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(-1), Inches(-1), Inches(5), Inches(5),
            self.colors["accent2"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(9), Inches(3), Inches(6), Inches(6),
            self.colors["accent3"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(10), Inches(4), Inches(4), Inches(4),
            self.colors["accent4"], MSO_SHAPE.OVAL)

        self._add_rounded_rect(slide, Inches(2.5), Inches(1.8),
            Inches(8.3), Inches(3.5), self.colors["card_bg"])

        self._add_text(slide, Inches(2.5), Inches(2.0), Inches(8.3), Inches(1),
            icon, Pt(48), self.colors["accent3"], alignment=PP_ALIGN.CENTER)

        self._add_text(slide, Inches(2.5), Inches(3.0), Inches(8.3), Inches(1.2),
            title, Pt(36), self.colors["title"],
            bold=True, alignment=PP_ALIGN.CENTER)

        self._add_shape(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.06),
            self.colors["highlight"])

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_closing_slide(self, data: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide)

        title = data.get("title", "Thank You")
        bullets = data.get("bullet_points", [])

        self._add_shape(slide, Inches(10), Inches(-1), Inches(5), Inches(5),
            self.colors["accent1"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(11), Inches(0), Inches(3.5), Inches(3.5),
            self.colors["accent2"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(-2), Inches(4), Inches(4), Inches(4),
            self.colors["accent3"], MSO_SHAPE.OVAL)
        self._add_shape(slide, Inches(7), Inches(5), Inches(3), Inches(3),
            self.colors["accent4"], MSO_SHAPE.OVAL)

        self._add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "🎯", Pt(60), self.colors["accent1"], alignment=PP_ALIGN.CENTER)

        self._add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
            title, Pt(44), self.colors["title"],
            bold=True, alignment=PP_ALIGN.CENTER)

        self._add_shape(slide, Inches(5), Inches(4.2), Inches(3.3), Inches(0.06),
            self.colors["highlight"])

        if bullets:
            for i, bullet in enumerate(bullets):
                self._add_text(slide,
                    Inches(2), Inches(4.6 + i * 0.5), Inches(9), Inches(0.5),
                    f"✦  {bullet}", Pt(16), self.colors["text"],
                    alignment=PP_ALIGN.CENTER)

        self._add_rounded_rect(slide, Inches(4), Inches(6), Inches(5.3), Inches(1),
            self.colors["card_bg"])
        self._add_text(slide, Inches(4.2), Inches(6.1), Inches(5), Inches(0.4),
            "Generated with ✦ newd2p AI", Pt(13), self.colors["muted"],
            alignment=PP_ALIGN.CENTER)
        self._add_text(slide, Inches(4.2), Inches(6.5), Inches(5), Inches(0.4),
            "Document → Narrative → Presentation", Pt(10), self.colors["muted"],
            alignment=PP_ALIGN.CENTER)

        notes = data.get("speaker_notes", "")
        if self.include_speaker_notes and notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _extract_json(self, text: str) -> dict:
        try:
            start = text.find("{")
            end = text.rfind("}") + 1
            if start != -1 and end > start:
                return json.loads(text[start:end])
        except:
            pass
        return {"title": "Presentation", "slides": []}
